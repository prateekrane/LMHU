from pptx import Presentation
from pptx.util import Inches
import logging

class PPTBuilder:
    def __init__(self):
        self.prs = Presentation()

    def run_command(self, command):
        action = command.get("action")
        title = command.get("title", "")
        content = command.get("content", "")
        try:
            if action == "add_slide":
                slide_layout = self.prs.slide_layouts[1]
                slide = self.prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = title
                slide.placeholders[1].text = content
            else:
                logging.warning(f"Unknown PPT command: {action}")
        except Exception as e:
            logging.error(f"Failed to run PPT command: {command}", exc_info=True)

    def save(self, filepath):
        self.prs.save(filepath)