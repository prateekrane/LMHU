from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.chart.data import CategoryChartData
import os
from datetime import datetime
import comtypes
import comtypes.client
import tempfile
import win32com.client

def _get_theme_colors(theme_name):
    # Define theme-specific color palettes
    color_palettes = {
        "Professional": {
            "primary": RGBColor(0, 112, 192),    # Professional blue
            "secondary": RGBColor(0, 176, 240),  # Light blue
            "accent": RGBColor(255, 255, 255),   # White
            "text": RGBColor(68, 68, 68)         # Dark gray
        },
        "Creative": {
            "primary": RGBColor(255, 89, 94),    # Vibrant red
            "secondary": RGBColor(255, 202, 58), # Bright yellow
            "accent": RGBColor(138, 201, 38),    # Fresh green
            "text": RGBColor(25, 25, 25)         # Near black
        },
        "Corporate": {
            "primary": RGBColor(31, 73, 125),    # Deep blue
            "secondary": RGBColor(79, 129, 189), # Medium blue
            "accent": RGBColor(192, 80, 77),     # Corporate red
            "text": RGBColor(89, 89, 89)         # Medium gray
        },
        "Modern": {
            "primary": RGBColor(45, 55, 72),     # Dark slate
            "secondary": RGBColor(237, 242, 247),# Light gray
            "accent": RGBColor(66, 153, 225),    # Modern blue
            "text": RGBColor(26, 32, 44)         # Dark slate
        },
        "Elegant": {
            "primary": RGBColor(72, 52, 52),     # Deep burgundy
            "secondary": RGBColor(192, 152, 152),# Soft rose
            "accent": RGBColor(240, 230, 220),   # Cream
            "text": RGBColor(51, 51, 51)         # Rich black
        }
    }
    return color_palettes.get(theme_name, color_palettes["Professional"])

def _apply_theme(prs, theme_name):
    # Enhanced theme configurations with specific requirements
    theme_map = {
        "Professional": {
            "layout": "Office Theme",
            "elements": ["transition"],
            "title_font_size": 40,
            "subtitle_font_size": 28,
            "content_font_size": 20,
            "bullet_font_size": 18,
            "transition": {
                "effect": "fade",
                "duration": 0.5
            }
        },
        "Creative": {
            "layout": "Facet",
            "elements": ["transition", "animation"],
            "title_font_size": 44,
            "subtitle_font_size": 32,
            "content_font_size": 24,
            "bullet_font_size": 22,
            "animation_settings": {
                "title": "zoom",
                "content": "float"
            },
            "transition": {
                "effect": "wipe",
                "duration": 0.3
            }
        },
        "Corporate": {
            "layout": "Office Theme",
            "elements": ["transition", "animation"],
            "title_font_size": 36,
            "subtitle_font_size": 28,
            "content_font_size": 20,
            "bullet_font_size": 18,
            "animation_settings": {
                "title": "fly_in",
                "content": "fade"
            },
            "transition": {
                "effect": "push",
                "duration": 0.5
            }
        },
        "Modern": {
            "layout": "Ion",
            "elements": ["transition", "animation"],
            "title_font_size": 42,
            "subtitle_font_size": 30,
            "content_font_size": 22,
            "bullet_font_size": 20,
            "animation_settings": {
                "title": "float",
                "content": "zoom"
            },
            "transition": {
                "effect": "cut",
                "duration": 0.3
            }
        },
        "Elegant": {
            "layout": "Office Theme",
            "elements": ["transition", "animation"],
            "title_font_size": 38,
            "subtitle_font_size": 28,
            "content_font_size": 20,
            "bullet_font_size": 18,
            "animation_settings": {
                "title": "fade",
                "content": "fly_in"
            },
            "transition": {
                "effect": "dissolve",
                "duration": 0.7
            }
        }
    }
    return theme_map.get(theme_name, theme_map["Professional"])

def _apply_color_scheme(prs, color_scheme, theme_name):
    # Get theme-specific colors
    theme_colors = _get_theme_colors(theme_name)
    
    # Override with color scheme if specified
    color_schemes = {
        "Default": theme_colors,
        "Blue": {
            "primary": RGBColor(0, 112, 192),
            "secondary": RGBColor(0, 176, 240),
            "accent": RGBColor(255, 255, 255),
            "text": RGBColor(68, 68, 68)
        },
        "Green": {
            "primary": RGBColor(0, 176, 80),
            "secondary": RGBColor(146, 208, 80),
            "accent": RGBColor(255, 255, 255),
            "text": RGBColor(68, 68, 68)
        },
        "Red": {
            "primary": RGBColor(192, 0, 0),
            "secondary": RGBColor(255, 0, 0),
            "accent": RGBColor(255, 255, 255),
            "text": RGBColor(68, 68, 68)
        },
        "Purple": {
            "primary": RGBColor(112, 48, 160),
            "secondary": RGBColor(149, 55, 53),
            "accent": RGBColor(255, 255, 255),
            "text": RGBColor(68, 68, 68)
        },
        "Orange": {
            "primary": RGBColor(255, 140, 0),
            "secondary": RGBColor(255, 192, 0),
            "accent": RGBColor(255, 255, 255),
            "text": RGBColor(68, 68, 68)
        }
    }
    return color_schemes.get(color_scheme, theme_colors)

def _apply_font(shape, font_name, font_size=None, color=None, bold=False, italic=False):
    if hasattr(shape, "text_frame"):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                if font_size:
                    run.font.size = Pt(font_size)
                if color:
                    run.font.color.rgb = color
                if bold:
                    run.font.bold = True
                if italic:
                    run.font.italic = True

def _add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    return shape

def _add_chart(slide, chart_type, left, top, width, height, data, theme_colors):
    # Create chart data object
    chart_data = CategoryChartData()
    
    # Add categories
    chart_data.categories = data.get("categories", ["Q1", "Q2", "Q3", "Q4"])
    
    # Add series
    for series_name, series_values in data.get("series", []):
        chart_data.add_series(series_name, series_values)
    
    # Add chart to slide
    chart = slide.shapes.add_chart(
        chart_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data
    ).chart
    
    # Apply theme colors to chart
    chart.chart_style = 2  # Use a clean style
    
    # Format chart elements
    if hasattr(chart, 'has_legend') and chart.has_legend:
        chart.legend.font.size = Pt(10)
        chart.legend.font.color.rgb = theme_colors["text"]
    
    if hasattr(chart, 'chart_title') and chart.chart_title:
        chart.chart_title.font.size = Pt(14)
        chart.chart_title.font.color.rgb = theme_colors["primary"]
    
    # Format series
    for idx, series in enumerate(chart.series):
        series.format.fill.solid()
        # Alternate between primary and secondary colors
        color = theme_colors["primary"] if idx % 2 == 0 else theme_colors["secondary"]
        series.format.fill.fore_color.rgb = color
    
    return chart

def _get_topic_based_elements(topic, theme_info, theme_colors):
    """Determine appropriate shapes and elements based on the presentation topic."""
    topic = topic.lower()
    
    # Define topic categories and their associated elements using only basic shapes
    topic_elements = {
        "business": {
            "shapes": [
                (MSO_SHAPE.RECTANGLE, "Business growth chart"),
                (MSO_SHAPE.FLOWCHART_PROCESS, "Business process"),
                (MSO_SHAPE.CUBE, "Business building blocks"),
                (MSO_SHAPE.ROUNDED_RECTANGLE, "Business card style")
            ],
            "positions": [
                (7, 1, 2, 1),  # right side
                (0.5, 0.5, 1.5, 0.5),  # top left
                (8, 2, 1, 1),  # bottom right
                (0, 5, 2, 0.3)  # bottom left
            ]
        },
        "technology": {
            "shapes": [
                (MSO_SHAPE.HEXAGON, "Tech network"),
                (MSO_SHAPE.ROUNDED_RECTANGLE, "Technology cycle"),
                (MSO_SHAPE.ROUNDED_RECTANGLE, "Tech interface"),
                (MSO_SHAPE.FLOWCHART_DECISION, "Tech decision point")
            ],
            "positions": [
                (7.5, 1, 1.5, 1.5),
                (0.5, 0.5, 1.2, 1.2),
                (8, 2, 1.5, 0.8),
                (0.5, 5, 1.5, 0.8)
            ]
        },
        "education": {
            "shapes": [
                (MSO_SHAPE.OVAL, "Learning cycle"),
                (MSO_SHAPE.ROUNDED_RECTANGLE, "Achievement"),
                (MSO_SHAPE.ROUNDED_RECTANGLE, "Book style"),
                (MSO_SHAPE.FLOWCHART_PROCESS, "Learning process")
            ],
            "positions": [
                (7, 1, 1.5, 1),
                (0.5, 0.5, 1, 1),
                (8, 2, 1.2, 1.5),
                (0.5, 5, 1.5, 0.5)
            ]
        },
        "health": {
            "shapes": [
                (MSO_SHAPE.OVAL, "Health symbol"),
                (MSO_SHAPE.OVAL, "Medical cycle"),
                (MSO_SHAPE.ROUNDED_RECTANGLE, "Medical card"),
                (MSO_SHAPE.FLOWCHART_PROCESS, "Health process")
            ],
            "positions": [
                (7.5, 1, 1, 1),
                (0.5, 0.5, 1.2, 1.2),
                (8, 2, 1.2, 0.8),
                (0.5, 5, 1.5, 0.5)
            ]
        },
        "environment": {
            "shapes": [
                (MSO_SHAPE.OVAL, "Sun energy"),
                (MSO_SHAPE.ROUNDED_RECTANGLE, "Cloud/atmosphere"),
                (MSO_SHAPE.ROUNDED_RECTANGLE, "Water/waves"),
                (MSO_SHAPE.OVAL, "Earth cycle")
            ],
            "positions": [
                (7.5, 0.5, 1, 1),
                (0.5, 0.5, 1.2, 0.8),
                (8, 2, 1.5, 0.5),
                (0.5, 5, 1.2, 1.2)
            ]
        },
        "finance": {
            "shapes": [
                (MSO_SHAPE.RECTANGLE, "Financial chart"),
                (MSO_SHAPE.ROUNDED_RECTANGLE, "Financial cycle"),
                (MSO_SHAPE.CUBE, "Investment blocks"),
                (MSO_SHAPE.ROUNDED_RECTANGLE, "Financial card")
            ],
            "positions": [
                (7, 1, 2, 1),
                (0.5, 0.5, 1.2, 1.2),
                (8, 2, 1, 1),
                (0.5, 5, 1.5, 0.3)
            ]
        }
    }
    
    # Default elements if no specific topic is matched
    default_elements = {
        "shapes": [
            (MSO_SHAPE.ROUNDED_RECTANGLE, "Decorative element"),
            (MSO_SHAPE.OVAL, "Design element"),
            (MSO_SHAPE.FLOWCHART_PROCESS, "Process element"),
            (MSO_SHAPE.ROUNDED_RECTANGLE, "Cycle element")
        ],
        "positions": [
            (7, 1, 2, 0.5),
            (0.5, 0.5, 1.2, 1.2),
            (8, 2, 1.5, 0.5),
            (0.5, 5, 1.5, 0.5)
        ]
    }
    
    # Determine the most relevant topic category
    matched_topic = None
    for category in topic_elements:
        if category in topic:
            matched_topic = category
            break
    
    # Get elements based on matched topic or use default
    elements = topic_elements.get(matched_topic, default_elements)
    
    return elements["shapes"], elements["positions"]

def _apply_theme_elements(slide, theme_info, formatting, theme_colors, topic=""):
    if not formatting or "elements" not in formatting:
        return

    for element in formatting.get("elements", []):
        element_type = element.get("type", "")
        
        if element_type == "shape" and "shape" in theme_info["elements"]:
            # Get topic-based elements
            shapes, positions = _get_topic_based_elements(topic, theme_info, theme_colors)
            
            # Apply shapes based on theme and topic
            for (shape_type, _), (left, top, width, height) in zip(shapes, positions):
                # Adjust colors based on theme
                if theme_info["layout"] == "Facet":  # Creative theme
                    fill_color = theme_colors["accent"] if shape_type in [MSO_SHAPE.SUN, MSO_SHAPE.CLOUD] else theme_colors["primary"]
                    line_color = theme_colors["primary"] if shape_type in [MSO_SHAPE.STAR_8_POINT, MSO_SHAPE.HEART] else None
                elif theme_info["layout"] == "Ion":  # Modern theme
                    fill_color = theme_colors["secondary"] if shape_type in [MSO_SHAPE.CIRCULAR_ARROW, MSO_SHAPE.WAVE] else theme_colors["primary"]
                    line_color = None
                elif "Elegant" in theme_info["layout"]:  # Elegant theme
                    fill_color = theme_colors["accent"] if shape_type in [MSO_SHAPE.STAR_12_POINT, MSO_SHAPE.OVAL] else theme_colors["primary"]
                    line_color = theme_colors["primary"] if shape_type in [MSO_SHAPE.STAR_12_POINT, MSO_SHAPE.CIRCULAR_ARROW] else None
                else:  # Professional and Corporate themes
                    fill_color = theme_colors["primary"]
                    line_color = None
                
                _add_shape(slide, shape_type, left, top, width, height,
                          fill_color=fill_color,
                          line_color=line_color)
        
        elif element_type == "chart" and "chart" in theme_info["elements"]:
            # Create sample chart data based on topic
            chart_data = _get_topic_based_chart_data(topic)
            chart_type = theme_info["chart_types"][0]
            _add_chart(slide, chart_type, 5, 2, 4, 3, chart_data, theme_colors)
        
        elif element_type == "design" and "design" in theme_info["elements"]:
            # Add design elements based on theme and topic
            if theme_info["layout"] == "Ion":  # Modern theme
                _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0, 0, 10, 0.5,
                          fill_color=theme_colors["primary"])
                _add_shape(slide, MSO_SHAPE.OVAL, 7, 1, 2, 0.3,
                          fill_color=theme_colors["secondary"])
            elif "Elegant" in theme_info["layout"]:  # Elegant theme
                _add_shape(slide, MSO_SHAPE.OVAL, 7, 1, 1, 1,
                          fill_color=theme_colors["accent"])
                _add_shape(slide, MSO_SHAPE.OVAL, 8, 0, 1, 1,
                          fill_color=theme_colors["secondary"])

def _get_topic_based_chart_data(topic):
    """Generate appropriate chart data based on the presentation topic."""
    pass

def _apply_comtypes_transitions(filepath, transition_effect, transition_duration):
    """Apply slide transitions using comtypes based on the provided effect and duration."""
    comtypes.CoInitialize()  # Initialize COM
    try:
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1

        presentation = powerpoint.Presentations.Open(filepath)

        for slide in presentation.Slides:
            slide.SlideShowTransition.EntryEffect = getattr(
                comtypes.gen.PowerPoint.PpEntryEffect,
                f"ppEffect{transition_effect.capitalize()}",
                0
            )
            slide.SlideShowTransition.Duration = transition_duration

        presentation.Save()
        presentation.Close()
        powerpoint.Quit()
    finally:
        comtypes.CoUninitialize()  # Uninitialize COM

def _apply_text_box_animation_with_pywin(filepath, slide_index, shape_index, animation_type):
    """Apply animation to a text box using pywin32."""
    powerpoint = None
    presentation = None
    try:
        # Create PowerPoint application instance
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.DisplayAlerts = False  # Disable alerts
        powerpoint.Visible = False  # Make PowerPoint invisible

        # Open presentation
        presentation = powerpoint.Presentations.Open(
            filepath,
            WithWindow=False,  # Don't show window
            ReadOnly=False     # Read-write mode
        )

        # Get the slide and shape
        slide = presentation.Slides(slide_index + 1)
        shape = slide.Shapes(shape_index + 1)

        # Define animation effects mapping
        animation_effects = {
            "fade": 0x17,      # ppEffectFade
            "float_in": 0x0B,  # ppEffectFloat
            "zoom": 0x0E,      # ppEffectZoom
            "fly_in": 0x04     # ppEffectFly
        }

        # Get effect number (default to fade)
        effect_number = animation_effects.get(animation_type.lower(), 0x17)

        # Apply animation settings
        animation = shape.AnimationSettings
        animation.Animate = True  # Enable animation
        animation.EntryEffect = effect_number

        # Save the presentation
        presentation.Save()

    except Exception as e:
        print(f"Animation error: {str(e)}")
    finally:
        if presentation:
            try:
                presentation.Close()
            except Exception as e:
                print(f"Error closing presentation: {str(e)}")
        if powerpoint:
            try:
                powerpoint.Quit()
            except Exception as e:
                print(f"Error quitting PowerPoint: {str(e)}")

def generate_ppt_doc(data):
    temp_files = []
    powerpoint = None
    try:
        prs = Presentation()

        # Get theme and color settings
        theme_info = _apply_theme(prs, data.get("theme", "Professional"))
        theme_colors = _apply_color_scheme(prs, data.get("color_scheme", "Default"), data.get("theme", "Professional"))
        font_name = data.get("font", "Calibri")
        theme_name = data.get("theme", "Professional")

        # Create title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        # Apply title and subtitle
        title_shape = slide.shapes.title
        title_shape.text = data.get("title", "Untitled Presentation")
        _apply_font(title_shape, font_name, theme_info["title_font_size"], theme_colors["primary"], bold=True)

        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = data.get("subtitle", "")
        _apply_font(subtitle_shape, font_name, theme_info["subtitle_font_size"], theme_colors["secondary"])

        # Process remaining slides
        slides_data = data.get("slides", [])
        include_toc = data.get("include_toc", True)

        # Create all slides first
        for idx, slide_data in enumerate(slides_data):
            is_toc = slide_data.get("title", "").lower() == "table of contents"
            layout = prs.slide_layouts[2] if is_toc else prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)

            # Apply content
            title_shape = slide.shapes.title
            title_shape.text = slide_data.get("title", "")
            _apply_font(title_shape, font_name, theme_info["title_font_size"], theme_colors["primary"], bold=True)

            if not is_toc:
                content_shape = slide.placeholders[1]
                content_shape.text = slide_data.get("content", "")
                _apply_font(content_shape, font_name, theme_info["content_font_size"], theme_colors["text"])
            else:
                # Handle TOC slide
                left, top, width, height = Inches(1.5), Inches(2.2), Inches(7), Inches(4)
                toc_box = slide.shapes.add_textbox(left, top, width, height)
                tf = toc_box.text_frame
                tf.word_wrap = True

                content = slide_data.get("content", "")
                for line in content.split('\n'):
                    if line.strip():
                        p = tf.add_paragraph()
                        p.text = line.strip()
                        p.alignment = PP_ALIGN.LEFT
                        for run in p.runs:
                            run.font.name = font_name
                            run.font.size = Pt(theme_info["content_font_size"])
                            run.font.color.rgb = theme_colors["text"]

        # Generate final filename
        downloads_path = os.path.expanduser("~/Downloads")
        os.makedirs(downloads_path, exist_ok=True)

        title = data.get("title", "Untitled").strip()
        clean_title = "".join(c for c in title if c.isalnum() or c.isspace())
        clean_title = clean_title.replace(" ", "_")[:30]
        timestamp = datetime.now().strftime("%Y%m%d_%H%M")
        filename = f"{clean_title}_{timestamp}.pptx"
        final_filepath = os.path.join(downloads_path, filename)

        # Save to final location
        prs.save(final_filepath)

        # Apply animations only to the final file
        if "animation" in theme_info["elements"]:
            try:
                comtypes.CoInitialize()
                powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
                powerpoint.DisplayAlerts = 0
                powerpoint.Visible = 0

                # Open the final presentation
                presentation = powerpoint.Presentations.Open(
                    final_filepath,
                    WithWindow=0,
                    ReadOnly=0
                )

                theme_animations = {
                    "Creative": {"title": "zoom", "content": "float"},
                    "Corporate": {"title": "fly_in", "content": "fade"},
                    "Modern": {"title": "float", "content": "zoom"},
                    "Elegant": {"title": "fade", "content": "fly_in"}
                }

                animations = theme_animations.get(theme_name, {})
                if animations:
                    for idx, slide in enumerate(presentation.Slides):
                        if idx == 0:  # Title slide
                            shape = slide.Shapes.Title
                            _apply_text_box_animation_with_pywin(final_filepath, idx, shape.ZOrderPosition - 1, animations.get("title", "fade"))
                        else:  # Content slides
                            for shape in slide.Shapes:
                                if shape.HasTextFrame:
                                    _apply_text_box_animation_with_pywin(final_filepath, idx, shape.ZOrderPosition - 1, animations.get("content", "fade"))

                presentation.Save()
                presentation.Close()

            except Exception as e:
                print(f"Animation error: {str(e)}")
            finally:
                if powerpoint:
                    powerpoint.Quit()
                comtypes.CoUninitialize()

        return final_filepath

    except Exception as e:
        print(f"Error generating presentation: {str(e)}")
        raise
    finally:
        for temp_file in temp_files:
            try:
                os.remove(temp_file)
            except Exception as e:
                print(f"Error cleaning up temp file: {str(e)}")